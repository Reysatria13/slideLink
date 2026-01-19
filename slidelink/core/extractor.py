"""
SlideLink Extractor - PPTX text extraction with bounding box metadata

Extracts all text from a PowerPoint file along with position and size
information needed for constraint-aware translation.
"""

from pptx import Presentation
from pptx.util import Emu
from typing import List, Optional
from dataclasses import dataclass, field


@dataclass
class BoundingBox:
    """Bounding box for a shape in pixels (96 DPI)"""
    left: int
    top: int
    width: int
    height: int

    # EMU values for reference
    left_emu: int = 0
    top_emu: int = 0
    width_emu: int = 0
    height_emu: int = 0


@dataclass
class Paragraph:
    """A paragraph within a text shape"""
    text: str
    font_size_pt: Optional[float] = None
    font_bold: bool = False
    alignment: Optional[str] = None


@dataclass
class ShapeData:
    """Extracted data for a single shape"""
    shape_id: int
    shape_name: str
    slide_index: int
    text: str
    paragraphs: List[Paragraph]
    bounding_box: BoundingBox
    estimated_max_chars: int
    primary_font_size_pt: float = 16.0


@dataclass
class ExtractionResult:
    """Result of extracting a presentation"""
    source_file: str
    slide_width_px: int
    slide_height_px: int
    slide_count: int
    shapes: List[ShapeData] = field(default_factory=list)


def emu_to_px(emu: int, dpi: int = 96) -> int:
    """Convert EMU (English Metric Units) to pixels"""
    return int(emu / 914400 * dpi)


def emu_to_pt(emu: int) -> float:
    """Convert EMU to points"""
    return emu / 914400 * 72


def estimate_max_chars(width_px: int, height_px: int, font_size_pt: float = 16.0) -> int:
    """
    Estimate maximum characters that fit in a bounding box.

    This is approximate - actual fit depends on font family,
    character width variance, and line spacing.
    """
    char_width = font_size_pt * 0.6  # Average character width
    line_height = font_size_pt * 1.5  # Line height

    chars_per_line = int(width_px / char_width) if char_width > 0 else 50
    num_lines = int(height_px / line_height) if line_height > 0 else 3

    return chars_per_line * num_lines


def extract_presentation(pptx_path: str) -> ExtractionResult:
    """
    Extract all text with bounding box metadata from a PPTX file.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        ExtractionResult with all shapes and their metadata
    """
    prs = Presentation(pptx_path)

    result = ExtractionResult(
        source_file=pptx_path,
        slide_width_px=emu_to_px(prs.slide_width),
        slide_height_px=emu_to_px(prs.slide_height),
        slide_count=len(prs.slides),
        shapes=[]
    )

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Collect paragraphs
            paragraphs = []
            full_text_parts = []
            primary_font_size = 16.0

            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue

                full_text_parts.append(para_text)

                # Get font info from first run
                font_size = None
                font_bold = False
                if para.runs:
                    if para.runs[0].font.size:
                        font_size = para.runs[0].font.size.pt
                        primary_font_size = font_size
                    font_bold = para.runs[0].font.bold or False

                paragraphs.append(Paragraph(
                    text=para_text,
                    font_size_pt=font_size,
                    font_bold=font_bold,
                    alignment=str(para.alignment) if para.alignment else None
                ))

            if not paragraphs:
                continue

            full_text = "\n".join(full_text_parts)

            # Create bounding box
            bbox = BoundingBox(
                left=emu_to_px(shape.left),
                top=emu_to_px(shape.top),
                width=emu_to_px(shape.width),
                height=emu_to_px(shape.height),
                left_emu=shape.left,
                top_emu=shape.top,
                width_emu=shape.width,
                height_emu=shape.height
            )

            # Estimate max characters
            max_chars = estimate_max_chars(
                bbox.width,
                bbox.height,
                primary_font_size
            )

            shape_data = ShapeData(
                shape_id=shape.shape_id,
                shape_name=shape.name,
                slide_index=slide_idx,
                text=full_text,
                paragraphs=paragraphs,
                bounding_box=bbox,
                estimated_max_chars=max_chars,
                primary_font_size_pt=primary_font_size
            )

            result.shapes.append(shape_data)

    return result
