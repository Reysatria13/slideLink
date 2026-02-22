"""
SlideLink Reconstructor - PPTX reconstruction with translated text

Opens the original PPTX, replaces text in each shape while preserving
all formatting, positioning, and styling.
"""

import os
from typing import Dict, List
from pptx import Presentation

from .translator import TranslationResult


def reconstruct_presentation(
    original_path: str,
    translations: List[TranslationResult],
    output_path: str = None
) -> str:
    """
    Create a new PPTX with translated text while preserving layout.

    Args:
        original_path: Path to the original PPTX file
        translations: List of TranslationResult from translator
        output_path: Optional output path. If None, auto-generates from input.

    Returns:
        Path to the output file
    """
    # Auto-generate output path if not provided
    if output_path is None:
        base, ext = os.path.splitext(original_path)
        output_path = f"{base}_translated{ext}"

    # Open original as template
    prs = Presentation(original_path)

    # Create lookup: (slide_index, shape_id) -> translated_text
    translation_map = {
        (t.slide_index, t.shape_id): t.translated_text for t in translations
    }

    # Update text in each shape
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            key = (slide_idx, shape.shape_id)
            if key not in translation_map:
                continue

            translated_text = translation_map[key]

            # Split translated text by newlines to match paragraphs
            translated_parts = translated_text.split("\n")

            # Update each paragraph
            for i, para in enumerate(shape.text_frame.paragraphs):
                if not para.text.strip():
                    continue

                # Get the translated text for this paragraph
                if i < len(translated_parts):
                    new_text = translated_parts[i]
                else:
                    new_text = ""

                # Preserve formatting by updating runs
                if para.runs:
                    # Put all text in first run, clear others
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    # No runs, just set text directly
                    para.text = new_text

    # Save to output path
    prs.save(output_path)

    return output_path


def reconstruct_with_custom_translations(
    original_path: str,
    translations: dict,
    output_path: str = None
) -> str:
    """
    Reconstruct with a custom translation dictionary.

    This is useful for the API when user edits translations.

    Args:
        original_path: Path to the original PPTX file
        translations: Dict mapping (slide_index, shape_id) -> translated_text
        output_path: Optional output path

    Returns:
        Path to the output file
    """
    if output_path is None:
        base, ext = os.path.splitext(original_path)
        output_path = f"{base}_translated{ext}"

    prs = Presentation(original_path)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            key = (slide_idx, shape.shape_id)
            if key not in translations:
                continue

            translated_text = translations[key]
            translated_parts = translated_text.split("\n")

            for i, para in enumerate(shape.text_frame.paragraphs):
                if not para.text.strip():
                    continue

                if i < len(translated_parts):
                    new_text = translated_parts[i]
                else:
                    new_text = ""

                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.text = new_text

    prs.save(output_path)
    return output_path
