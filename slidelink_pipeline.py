"""
SlideLink Complete Pipeline Demo

This demonstrates the full workflow:
1. EXTRACT: Parse PPTX and get text with bounding box metadata
2. TRANSLATE: Send to AI with constraints (simulated here)
3. RECONSTRUCT: Build new PPTX with translated text

This is a proof-of-concept for the SlideLink system.
"""

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.text import PP_ALIGN
import json
import copy

# ============================================================
# CONFIGURATION
# ============================================================

# Simulated translations (in real system, this would call Gemini API)
MOCK_TRANSLATIONS = {
    "2025年度 事業計画": {
        "translation": "FY2025 Business Plan",
        "fits": True
    },
    "先進AI戦略企画部": {
        "translation": "Advanced AI Strategy Division",
        "fits": True
    },
    "2025年1月": {
        "translation": "January 2025",
        "fits": True
    },
    "主要な取り組み": {
        "translation": "Key Initiatives",
        "fits": True
    },
    "AI翻訳システムの開発": {
        "translation": "AI Translation System Development",
        "fits": True
    },
    "社内プレゼンテーションの自動翻訳を実現し、グローバルコミュニケーションを促進します。": {
        "translation": "Enabling automatic translation of internal presentations to promote global communication.",
        "fits": True
    },
    "レイアウト保持機能": {
        "translation": "Layout Preservation Feature",
        "fits": True
    },
    "元のスライドデザインを維持しながら、テキストのみを翻訳します。": {
        "translation": "Translates text while maintaining original slide design.",
        "fits": True
    },
    "期待される効果": {
        "translation": "Expected Benefits",
        "fits": True
    },
    "• 翻訳時間: 2-3時間 → 5分": {
        "translation": "• Translation time: 2-3 hours → 5 min",
        "fits": True
    },
    "• コスト削減: 年間約500万円": {
        "translation": "• Cost savings: ~¥5M annually",
        "fits": True
    },
    "• 品質向上: 一貫した用語使用": {
        "translation": "• Quality improvement: Consistent terminology",
        "fits": True
    },
    "機密 - 社内使用限定": {
        "translation": "Confidential - Internal Use Only",
        "fits": True
    }
}

# ============================================================
# STEP 1: EXTRACTION
# ============================================================

def emu_to_px(emu, dpi=96):
    return int(emu / 914400 * dpi)

def extract_presentation(pptx_path):
    """Extract all text with metadata from PPTX"""
    prs = Presentation(pptx_path)
    
    extraction = {
        "source_file": pptx_path,
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": slide_idx,
            "shapes": []
        }
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Collect paragraphs
            paragraphs = []
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    # Get formatting
                    font_size = None
                    font_bold = False
                    if para.runs:
                        if para.runs[0].font.size:
                            font_size = para.runs[0].font.size
                        font_bold = para.runs[0].font.bold or False
                    
                    paragraphs.append({
                        "text": para.text,
                        "font_size_emu": font_size,
                        "font_bold": font_bold,
                        "alignment": str(para.alignment) if para.alignment else None
                    })
            
            if not paragraphs:
                continue
            
            shape_data = {
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "left_emu": shape.left,
                "top_emu": shape.top,
                "width_emu": shape.width,
                "height_emu": shape.height,
                "paragraphs": paragraphs
            }
            
            slide_data["shapes"].append(shape_data)
        
        extraction["slides"].append(slide_data)
    
    return extraction

# ============================================================
# STEP 2: TRANSLATION (with constraints)
# ============================================================

def translate_text(japanese_text, box_width_px, box_height_px, font_size_pt=16):
    """
    Translate text with layout constraints.
    
    In production, this would call Gemini API with a prompt like:
    
    ```
    Translate to English. The text must fit in a {box_width}px × {box_height}px box
    at {font_size}pt. Original: "{japanese_text}"
    
    If too long, provide a shorter alternative.
    ```
    
    For this demo, we use mock translations.
    """
    
    # Check if we have a mock translation
    if japanese_text in MOCK_TRANSLATIONS:
        result = MOCK_TRANSLATIONS[japanese_text]
        return {
            "original": japanese_text,
            "translation": result["translation"],
            "fits_box": result["fits"],
            "method": "mock"
        }
    
    # Fallback: return original with [NEEDS TRANSLATION] marker
    return {
        "original": japanese_text,
        "translation": f"[TRANSLATE: {japanese_text}]",
        "fits_box": True,
        "method": "fallback"
    }


def translate_extraction(extraction):
    """Translate all extracted text"""
    translated = copy.deepcopy(extraction)
    
    translation_log = []
    
    for slide in translated["slides"]:
        for shape in slide["shapes"]:
            box_width = emu_to_px(shape["width_emu"])
            box_height = emu_to_px(shape["height_emu"])
            
            for para in shape["paragraphs"]:
                font_size = 16
                if para["font_size_emu"]:
                    font_size = para["font_size_emu"] / 914400 * 72
                
                result = translate_text(
                    para["text"],
                    box_width,
                    box_height,
                    font_size
                )
                
                # Store translation
                para["original_text"] = para["text"]
                para["text"] = result["translation"]
                
                translation_log.append({
                    "original": result["original"],
                    "translated": result["translation"],
                    "fits": result["fits_box"]
                })
    
    translated["translation_log"] = translation_log
    return translated

# ============================================================
# STEP 3: RECONSTRUCTION
# ============================================================

def reconstruct_presentation(original_path, translated_data, output_path):
    """
    Create a new PPTX with translated text while preserving layout.
    
    Key insight: We DON'T modify the original file.
    Instead, we create a copy and update only the text content,
    keeping all positioning, styling, and structure intact.
    """
    
    # Open original as template
    prs = Presentation(original_path)
    
    # Map shape_id to translated content
    translation_map = {}
    for slide_data in translated_data["slides"]:
        for shape_data in slide_data["shapes"]:
            translation_map[shape_data["shape_id"]] = shape_data["paragraphs"]
    
    # Update text in each shape
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if shape.shape_id not in translation_map:
                continue
            
            translated_paras = translation_map[shape.shape_id]
            
            # Match paragraphs and update text
            for i, para in enumerate(shape.text_frame.paragraphs):
                if i < len(translated_paras):
                    # Preserve formatting, update text
                    if para.runs:
                        # Update text in first run, clear others
                        para.runs[0].text = translated_paras[i]["text"]
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = translated_paras[i]["text"]
    
    # Save
    prs.save(output_path)
    return output_path

# ============================================================
# MAIN EXECUTION
# ============================================================

def main():
    print("=" * 70)
    print("SlideLink Pipeline Demo")
    print("=" * 70)
    
    source_file = "/home/claude/python_pptx_demo/sample_japanese.pptx"
    output_file = "/home/claude/python_pptx_demo/sample_english.pptx"
    
    # STEP 1: Extract
    print("\n📥 STEP 1: EXTRACTION")
    print("-" * 40)
    extraction = extract_presentation(source_file)
    print(f"   Source: {source_file}")
    print(f"   Slides: {len(extraction['slides'])}")
    
    total_shapes = sum(len(s['shapes']) for s in extraction['slides'])
    print(f"   Text boxes: {total_shapes}")
    
    # STEP 2: Translate
    print("\n🌐 STEP 2: TRANSLATION (with constraints)")
    print("-" * 40)
    translated = translate_extraction(extraction)
    
    print(f"   Translations performed: {len(translated['translation_log'])}")
    print("\n   Translation Log:")
    for i, item in enumerate(translated['translation_log'], 1):
        fits_icon = "✅" if item['fits'] else "⚠️"
        print(f"   {i}. {fits_icon} \"{item['original'][:30]}...\"")
        print(f"       → \"{item['translated'][:30]}...\"")
    
    # STEP 3: Reconstruct
    print("\n📤 STEP 3: RECONSTRUCTION")
    print("-" * 40)
    output = reconstruct_presentation(source_file, translated, output_file)
    print(f"   Output: {output}")
    
    # Summary
    print("\n" + "=" * 70)
    print("✅ PIPELINE COMPLETE")
    print("=" * 70)
    print(f"""
   Source (Japanese): {source_file}
   Output (English):  {output_file}
   
   The translated PPTX preserves:
   • Original slide dimensions
   • Text box positions and sizes
   • Font formatting (size, bold, alignment)
   • All other elements (shapes, images, etc.)
   
   This is the core of SlideLink's layout-preserving translation!
""")
    
    # Save extraction data for reference
    with open("/home/claude/python_pptx_demo/pipeline_data.json", "w", encoding="utf-8") as f:
        # Convert EMU values for JSON
        export_data = {
            "source": source_file,
            "output": output_file,
            "translations": translated["translation_log"]
        }
        json.dump(export_data, f, ensure_ascii=False, indent=2)
    
    print("   Pipeline data saved to: pipeline_data.json")

if __name__ == "__main__":
    main()
