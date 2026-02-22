from pptx import Presentation
prs = Presentation('sample_japanese.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        print(shape.name, shape.shape_type, getattr(shape, 'rotation', 0))
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    color = None
                    if r.font.color and hasattr(r.font.color, 'rgb') and r.font.color.rgb:
                        color = str(r.font.color.rgb)
                    print("  font color:", color)
