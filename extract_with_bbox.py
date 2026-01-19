"""
SlideLink Core: PPTX Text Extraction with Bounding Box Metadata

This script demonstrates how to extract text from a PowerPoint file
along with its position and dimension information - the foundation
for constraint-aware translation.

Key concepts:
- EMU (English Metric Units): PowerPoint's internal unit (914400 EMU = 1 inch)
- Bounding Box: The rectangle that contains a text box (left, top, width, height)
- This metadata allows AI to know "how much space" the translation has
"""

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

# Conversion helpers
def emu_to_inches(emu):
    """Convert EMU to inches (914400 EMU = 1 inch)"""
    return emu / 914400

def emu_to_px(emu, dpi=96):
    """Convert EMU to pixels at given DPI"""
    inches = emu / 914400
    return int(inches * dpi)

def emu_to_pt(emu):
    """Convert EMU to points (914400 EMU = 72 pt)"""
    return emu / 914400 * 72


def extract_slide_metadata(pptx_path):
    """
    Extract all text boxes with their bounding box metadata from a PPTX file.
    
    Returns a structure like:
    {
        "slide_width": 1280,
        "slide_height": 720,
        "slides": [
            {
                "slide_number": 1,
                "shapes": [
                    {
                        "shape_id": 2,
                        "shape_type": "TEXT_BOX",
                        "text": "2025年度 事業計画",
                        "bounding_box": {
                            "left": 48,
                            "top": 192,
                            "width": 1184,
                            "height": 144
                        },
                        "font_size_pt": 44,
                        "estimated_max_chars": 25
                    }
                ]
            }
        ]
    }
    """
    prs = Presentation(pptx_path)
    
    result = {
        "file": pptx_path,
        "slide_width_px": emu_to_px(prs.slide_width),
        "slide_height_px": emu_to_px(prs.slide_height),
        "slide_width_inches": round(emu_to_inches(prs.slide_width), 2),
        "slide_height_inches": round(emu_to_inches(prs.slide_height), 2),
        "slides": []
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "shapes": []
        }
        
        for shape in slide.shapes:
            # Check if shape has text
            if not shape.has_text_frame:
                continue
            
            # Get all text from the shape
            full_text = ""
            paragraphs_data = []
            
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    full_text += para_text + "\n"
                    
                    # Get font size from first run if available
                    font_size = None
                    if para.runs:
                        font_size = para.runs[0].font.size
                        if font_size:
                            font_size = font_size.pt
                    
                    paragraphs_data.append({
                        "text": para_text,
                        "font_size_pt": font_size
                    })
            
            full_text = full_text.strip()
            if not full_text:
                continue
            
            # Get bounding box
            bbox = {
                "left_px": emu_to_px(shape.left),
                "top_px": emu_to_px(shape.top),
                "width_px": emu_to_px(shape.width),
                "height_px": emu_to_px(shape.height),
                "left_inches": round(emu_to_inches(shape.left), 2),
                "top_inches": round(emu_to_inches(shape.top), 2),
                "width_inches": round(emu_to_inches(shape.width), 2),
                "height_inches": round(emu_to_inches(shape.height), 2),
            }
            
            # Estimate max characters that fit
            # Rough estimate: average character width ≈ font_size * 0.5
            # This is a simplification - real calculation would consider font metrics
            avg_font_size = 16  # default
            for p in paragraphs_data:
                if p["font_size_pt"]:
                    avg_font_size = p["font_size_pt"]
                    break
            
            char_width_px = avg_font_size * 0.6  # approximate
            line_height_px = avg_font_size * 1.5
            
            chars_per_line = int(bbox["width_px"] / char_width_px) if char_width_px > 0 else 50
            lines_available = int(bbox["height_px"] / line_height_px) if line_height_px > 0 else 3
            estimated_max_chars = chars_per_line * lines_available
            
            shape_data = {
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type).replace("MSO_SHAPE_TYPE.", ""),
                "text": full_text,
                "char_count": len(full_text),
                "bounding_box": bbox,
                "paragraphs": paragraphs_data,
                "estimated_max_chars": estimated_max_chars,
                "translation_constraints": {
                    "max_chars_recommended": int(estimated_max_chars * 0.9),  # 90% to be safe
                    "box_width_px": bbox["width_px"],
                    "box_height_px": bbox["height_px"],
                    "primary_font_size_pt": avg_font_size
                }
            }
            
            slide_data["shapes"].append(shape_data)
        
        result["slides"].append(slide_data)
    
    return result


def print_extraction_report(data):
    """Print a human-readable report of extracted data"""
    print("=" * 70)
    print(f"PPTX EXTRACTION REPORT: {data['file']}")
    print("=" * 70)
    print(f"Slide Dimensions: {data['slide_width_px']}px × {data['slide_height_px']}px")
    print(f"                  ({data['slide_width_inches']}\" × {data['slide_height_inches']}\")")
    print()
    
    for slide in data["slides"]:
        print(f"\n{'─' * 70}")
        print(f"SLIDE {slide['slide_number']}")
        print(f"{'─' * 70}")
        
        for i, shape in enumerate(slide["shapes"], 1):
            bbox = shape["bounding_box"]
            constraints = shape["translation_constraints"]
            
            print(f"\n  📦 Shape {i}: {shape['shape_name']}")
            print(f"  ├── Type: {shape['shape_type']}")
            print(f"  ├── Position: ({bbox['left_px']}px, {bbox['top_px']}px)")
            print(f"  ├── Size: {bbox['width_px']}px × {bbox['height_px']}px")
            print(f"  ├── Text ({shape['char_count']} chars):")
            
            # Print text with indentation
            for line in shape["text"].split("\n"):
                print(f"  │     \"{line}\"")
            
            print(f"  │")
            print(f"  └── 🎯 TRANSLATION CONSTRAINTS:")
            print(f"       ├── Max recommended chars: {constraints['max_chars_recommended']}")
            print(f"       ├── Box width: {constraints['box_width_px']}px")
            print(f"       ├── Box height: {constraints['box_height_px']}px")
            print(f"       └── Font size: {constraints['primary_font_size_pt']}pt")


def generate_translation_prompt(shape_data):
    """
    Generate a prompt for AI translation with layout constraints.
    This is the KEY INNOVATION of SlideLink - passing bounding box
    constraints to the AI during translation.
    """
    constraints = shape_data["translation_constraints"]
    
    prompt = f"""TASK: Translate Japanese to English with layout constraints

ORIGINAL TEXT:
{shape_data['text']}

LAYOUT CONSTRAINTS:
- Bounding box: {constraints['box_width_px']}px × {constraints['box_height_px']}px
- Font size: {constraints['primary_font_size_pt']}pt
- Maximum recommended characters: {constraints['max_chars_recommended']}
- Original character count: {shape_data['char_count']}

TRANSLATION RULES:
1. Translation MUST fit within the bounding box
2. Prefer concise business English
3. If direct translation is too long, provide a shortened alternative
4. Preserve tone (formal/informal)
5. Keep proper nouns (company names, product names) unchanged

OUTPUT FORMAT (JSON):
{{
  "translation": "Your translation here",
  "char_count": <number>,
  "fits_original_box": true/false,
  "confidence": 0.0-1.0,
  "alternative_short": "Shorter version if main translation might overflow"
}}"""
    
    return prompt


# Main execution
if __name__ == "__main__":
    pptx_path = "/home/claude/python_pptx_demo/sample_japanese.pptx"
    
    # Extract all metadata
    print("\n🔍 Extracting text and bounding box metadata...\n")
    data = extract_slide_metadata(pptx_path)
    
    # Print report
    print_extraction_report(data)
    
    # Save JSON for reference
    with open("/home/claude/python_pptx_demo/extracted_metadata.json", "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"\n\n✅ Metadata saved to: extracted_metadata.json")
    
    # Show example translation prompt for first text box
    print("\n" + "=" * 70)
    print("EXAMPLE: AI TRANSLATION PROMPT WITH CONSTRAINTS")
    print("=" * 70)
    
    if data["slides"] and data["slides"][0]["shapes"]:
        first_shape = data["slides"][0]["shapes"][0]
        prompt = generate_translation_prompt(first_shape)
        print(prompt)
